from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
presentation = Presentation()

# Add a slide with a title and content layout
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_placeholder = slide.shapes.title
title_placeholder.text = "Flowchart Diagram"

# Add the flowchart image to the slide
img_path = "result/"  # Replace with the path to the generated flowchart image
left = Inches(1)  # Adjust the position on the slide
top = Inches(1)
slide.shapes.add_picture(img_path, left, top, height=Inches(5))  # Adjust the size as needed

# Save the presentation
presentation.save("/path/to/save/flowchart_presentation.pptx")  # Replace with the desired save path